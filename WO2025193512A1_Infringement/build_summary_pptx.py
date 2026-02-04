"""
Build condensed 3-slide EoU summaries for WO/2025/193512A1 — 4 targets.
Slide 1: Title + patent info + target info
Slide 2: Claim element mapping table (all 7 elements)
Slide 3: Key evidence, sources, risk assessment
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import os

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))

PATENT = "WO/2025/193512A1"
TITLE = "SINGLE SHOT 3D MODELLING FROM 2D IMAGE"
ASSIGNEE = "Carnegie Mellon University"
INVENTOR = "Marios Savvides et al."

ELEMENTS = [
    "(a) obtaining a 2D image of an object",
    "(b) classifying the object in the image",
    "(c) segmenting the object from the image",
    "(d) dimensionally sampling the segmented image",
    "(e) extracting texture information",
    "(f) generating a 3D mesh model",
    "(g) rendering texture onto 3D mesh",
]

DEP_CLAIMS = [
    "Cl.5: Neural network generates mesh",
    "Cl.9: Multiple views generated",
    "Cl.10: User-input manipulation",
    "Cl.11: 2D capture of each view",
]


def hex_to_rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARK = hex_to_rgb("1a1a2e")
ACCENT = hex_to_rgb("0f3460")
HIGHLIGHT = hex_to_rgb("e94560")
GRAY = hex_to_rgb("666666")
LIGHT_BG = hex_to_rgb("f0f0f5")
GREEN = hex_to_rgb("2d6a4f")
YELLOW_D = hex_to_rgb("b5830a")
RED_D = hex_to_rgb("9b2226")


def set_cell_fill(cell, color):
    tcPr = cell._tc.get_or_add_tcPr()
    solidFill = tcPr.makeelement(qn('a:solidFill'), {})
    srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': '{:02X}{:02X}{:02X}'.format(color[0], color[1], color[2])})
    solidFill.append(srgb)
    tcPr.append(solidFill)


def set_cell_text(cell, text, size=8, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    cell.text_frame.word_wrap = True
    # Reduce margins
    cell.text_frame.margin_top = Pt(2)
    cell.text_frame.margin_bottom = Pt(2)
    cell.text_frame.margin_left = Pt(3)
    cell.text_frame.margin_right = Pt(3)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, w, h, text, size=11, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_multiline(slide, left, top, w, h, lines, size=9, color=BLACK, bold_first=False):
    """lines is list of strings. First line optionally bold."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        if bold_first and i == 0:
            p.font.bold = True
        p.space_after = Pt(2)
    return box


TARGETS = {
    "Google_Shopping": {
        "company": "Google (Alphabet)",
        "product": "Google Shopping 3D Product Viewer",
        "status": "Deployed at scale — shoes, furniture, electronics, apparel",
        "how": "1-3 product photos → Veo video diffusion model → 360° interactive spin",
        "strength": "HIGH",
        "mapping": [
            ("STRONG", "Accepts 1-3 product images as input"),
            ("LIKELY", "Handles diverse categories; training organized by product type"),
            ("STRONG", "Blog confirms 'removing unwanted backgrounds'"),
            ("STRONG", "Predicts 3D priors, camera positions from sparse images"),
            ("STRONG", "Veo captures 'light, material, texture, and geometry'"),
            ("STRONG", "NeRF/mesh/volumetric 3D representations generated"),
            ("STRONG", "Photorealistic 360° textured product views"),
        ],
        "dep_match": ["YES — neural net (Veo)", "YES — 360° spins", "YES — user rotates", "YES — frame capture"],
        "sources": [
            "research.google/blog/bringing-3d-shoppable-products-online-with-generative-ai/",
            "support.google.com/merchants/answer/13671720",
            "support.google.com/merchants/answer/13675100",
        ],
        "key_risk": "Google has extensive patent portfolio for counter-assertion / cross-license.",
        "key_strength": "Google's own research blog documents nearly every claim element publicly.",
    },
    "Meshy_AI": {
        "company": "Meshy Inc.",
        "product": "Meshy AI 3D Model Generator",
        "status": "Commercial SaaS — 1M+ users, production 3D asset tool",
        "how": "Single image → multi-view generation → 3D mesh + PBR textures → export",
        "strength": "HIGH",
        "mapping": [
            ("STRONG", "Single image upload is primary input"),
            ("LIKELY", "Multi-view generation requires object type understanding"),
            ("STRONG", "Object isolated from background for 3D generation"),
            ("STRONG", "Smart Remesh: adjustable 1k-300k triangle/quad counts"),
            ("STRONG", "Full PBR: Diffuse, Roughness, Metallic, Normal maps"),
            ("STRONG", "Explicit mesh export: FBX, GLB, OBJ, STL, USDZ, BLEND"),
            ("STRONG", "PBR textures rendered onto mesh; browser 3D viewer"),
        ],
        "dep_match": ["YES — AI/NN-based", "YES — any angle preview", "YES — user rotates", "YES — export views"],
        "sources": [
            "meshy.ai/",
            "Meshy product documentation & marketing",
        ],
        "key_risk": "Smaller company — less cross-license leverage but also less litigation budget.",
        "key_strength": "Cleanest claim match. Explicit mesh + PBR textures = no ambiguity on elements (f)/(g).",
    },
    "Apple_SHARP": {
        "company": "Apple Inc.",
        "product": "SHARP + Spatial Scenes (iOS 26) + Vision Pro",
        "status": "Open-source research (Dec 2025); Spatial Scenes commercialized iOS 26",
        "how": "Single photo → neural net forward pass (<1s) → 3D Gaussian splat → render",
        "strength": "MEDIUM-HIGH",
        "mapping": [
            ("STRONG", "Single photograph input; CLI: 'sharp predict -i image.png'"),
            ("LIKELY", "Zero-shot generalization across object/scene types"),
            ("LIKELY", "Depth Pro estimates depth → implicit object/background separation"),
            ("STRONG", "Predicts metric-scale 3D Gaussian positions, scales, orientations"),
            ("STRONG", "Appearance parameters encode color/texture from input"),
            ("RISK", "Outputs Gaussian splats (.ply), NOT traditional polygon mesh"),
            ("STRONG", "Photorealistic rendering at 100+ FPS; Vision Pro integration"),
        ],
        "dep_match": ["YES — single NN pass", "YES — novel views", "YES — user moves viewpoint", "YES — frame render"],
        "sources": [
            "github.com/apple/ml-sharp",
            "machinelearning.apple.com/research/sharp-monocular-view",
            "apple.github.io/ml-sharp/",
            "huggingface.co/apple/Sharp",
        ],
        "key_risk": "Gaussian splats vs. 'mesh model' — claim construction question on element (f). Apple has massive patent portfolio for cross-license.",
        "key_strength": "Fully documented pipeline (open-source + paper). Spatial Scenes in iOS 26 = commercial deployment.",
    },
    "Amazon": {
        "company": "Amazon.com, Inc.",
        "product": "Amazon 3D/AR (Seller App + View in 3D + AWS Pipeline)",
        "status": "Millions of product listings; mandatory 3D models since Dec 2023",
        "how": "Seller scanning / reference photos → 3D model (GLB/GLTF) → View in 3D / AR",
        "strength": "HIGH",
        "mapping": [
            ("STRONG", "Seller App captures images; 2-10 reference photos required"),
            ("STRONG", "Category-dependent eligibility; different 3D experiences per category"),
            ("STRONG", "AWS blog: 'image segmentation and semantic labelling' in pipeline"),
            ("STRONG", "'Accurate product dimensions' required; 3D geometry extracted"),
            ("STRONG", "Photorealistic textures preserved in 3D product views"),
            ("STRONG", "GLB/GLTF output = explicit mesh format (vertices/faces)"),
            ("STRONG", "View in 3D: rotate/zoom any angle; 2X conversion uplift"),
        ],
        "dep_match": ["YES — AI pipeline", "YES — 360° rotation", "YES — user rotates/zooms", "YES — AR view capture"],
        "sources": [
            "sell.amazon.com/tools/3d-ar",
            "aws.amazon.com/blogs/spatial/3d-gaussian-splatting...",
            "aws.amazon.com/blogs/spatial/from-2d-to-3d...",
        ],
        "key_risk": "Amazon has massive litigation resources. 3D scanning (multi-image) vs. single-shot may be distinguished.",
        "key_strength": "Mandatory 3D models since Dec 2023. GLB/GLTF = unambiguous mesh. AWS docs confirm segmentation.",
    },
}


def build_summary(key, t):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ============ SLIDE 1: Title + Overview ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    # Top bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.9))  # 1 = rectangle
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK
    bar.line.fill.background()

    add_text(slide, 0.4, 0.1, 8, 0.4, f"Evidence of Use — {PATENT}", size=18, bold=True, color=WHITE)
    add_text(slide, 0.4, 0.5, 8, 0.3, TITLE, size=10, color=RGBColor(180, 180, 200))

    # Strength badge
    strength_color = GREEN if t["strength"] == "HIGH" else YELLOW_D if "MEDIUM" in t["strength"] else RED_D
    add_text(slide, 10.5, 0.15, 2.5, 0.6, f'INFRINGEMENT\n{t["strength"]}',
             size=14, bold=True, color=strength_color, align=PP_ALIGN.CENTER)

    # Target info block
    add_text(slide, 0.4, 1.2, 6, 0.35, t["company"], size=20, bold=True, color=DARK)
    add_text(slide, 0.4, 1.65, 8, 0.3, t["product"], size=12, color=ACCENT)
    add_text(slide, 0.4, 2.0, 8, 0.25, t["status"], size=9, color=GRAY)

    # How it works
    add_text(slide, 0.4, 2.6, 2, 0.3, "PIPELINE", size=9, bold=True, color=ACCENT)
    add_text(slide, 0.4, 2.9, 9, 0.3, t["how"], size=11, color=BLACK)

    # Patent info block
    add_text(slide, 0.4, 3.5, 2, 0.3, "PATENT", size=9, bold=True, color=ACCENT)
    info_lines = [
        f"Patent: {PATENT}",
        f"Assignee: {ASSIGNEE}",
        f"Inventor: {INVENTOR}",
        f"Claims Mapped: 1 (independent) + 5, 9, 10, 11 (dependent)",
    ]
    add_multiline(slide, 0.4, 3.8, 5, 1.0, info_lines, size=9, color=DARK)

    # Sources
    add_text(slide, 6.5, 3.5, 2, 0.3, "EVIDENCE SOURCES", size=9, bold=True, color=ACCENT)
    add_multiline(slide, 6.5, 3.8, 6.5, 1.2, t["sources"], size=8, color=RGBColor(0, 80, 160))

    # Risk / Strength summary boxes
    # Strength box
    sb = slide.shapes.add_shape(1, Inches(0.4), Inches(5.2), Inches(5.8), Inches(1.0))
    sb.fill.solid()
    sb.fill.fore_color.rgb = hex_to_rgb("e8f5e9")
    sb.line.fill.background()
    add_text(slide, 0.5, 5.1, 1.5, 0.3, "KEY STRENGTH", size=8, bold=True, color=GREEN)
    add_text(slide, 0.5, 5.4, 5.5, 0.7, t["key_strength"], size=9, color=DARK)

    # Risk box
    rb = slide.shapes.add_shape(1, Inches(6.5), Inches(5.2), Inches(6.5), Inches(1.0))
    rb.fill.solid()
    rb.fill.fore_color.rgb = hex_to_rgb("fff3e0")
    rb.line.fill.background()
    add_text(slide, 6.6, 5.1, 1.5, 0.3, "KEY RISK", size=8, bold=True, color=YELLOW_D)
    add_text(slide, 6.6, 5.4, 6.2, 0.7, t["key_risk"], size=9, color=DARK)

    # Footer
    add_text(slide, 0.4, 6.9, 4, 0.3, "CONFIDENTIAL — Attorney Work Product", size=7, bold=True, color=HIGHLIGHT)
    add_text(slide, 10, 6.9, 3, 0.3, "2/2/26  |  Slide 1 of 3", size=7, color=GRAY, align=PP_ALIGN.RIGHT)

    # ============ SLIDE 2: Claim Mapping Table ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK
    bar.line.fill.background()
    add_text(slide, 0.4, 0.1, 10, 0.35, f"Claim 1 Element Mapping — {t['company']}", size=14, bold=True, color=WHITE)

    # Main table: 8 rows (header + 7 elements), 3 columns
    rows, cols = 8, 3
    tbl = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(0.9), Inches(12.7), Inches(4.2)).table
    tbl.columns[0].width = Inches(3.2)
    tbl.columns[1].width = Inches(1.3)
    tbl.columns[2].width = Inches(8.2)

    # Header
    headers = ["Claim Element", "Match", "Evidence Summary"]
    header_cells = tbl.rows[0].cells
    for j, h in enumerate(headers):
        set_cell_text(header_cells[j], h, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        set_cell_fill(header_cells[j], (26, 26, 46))

    for i, (element, (match, evidence)) in enumerate(zip(ELEMENTS, t["mapping"])):
        row_cells = tbl.rows[i + 1].cells
        set_cell_text(row_cells[0], element, size=8, bold=True)
        match_color = GREEN if match == "STRONG" else YELLOW_D if match == "LIKELY" else RED_D
        set_cell_text(row_cells[1], match, size=8, bold=True, color=match_color, align=PP_ALIGN.CENTER)
        set_cell_text(row_cells[2], evidence, size=8)
        bg = (240, 240, 245) if i % 2 == 0 else (255, 255, 255)
        for j in range(3):
            set_cell_fill(row_cells[j], bg)

    # Dependent claims mini-table
    add_text(slide, 0.3, 5.3, 4, 0.3, "Dependent Claims", size=10, bold=True, color=ACCENT)
    rows2, cols2 = 2, 4
    tbl2 = slide.shapes.add_table(rows2, cols2, Inches(0.3), Inches(5.7), Inches(12.7), Inches(0.7)).table
    for j in range(4):
        tbl2.columns[j].width = Inches(3.175)

    for j, dc in enumerate(DEP_CLAIMS):
        set_cell_text(tbl2.rows[0].cells[j], dc, size=8, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        set_cell_fill(tbl2.rows[0].cells[j], (230, 230, 240))
        dm = t["dep_match"][j]
        dm_color = GREEN if dm.startswith("YES") else YELLOW_D
        set_cell_text(tbl2.rows[1].cells[j], dm, size=8, bold=True, color=dm_color, align=PP_ALIGN.CENTER)
        set_cell_fill(tbl2.rows[1].cells[j], (255, 255, 255))

    add_text(slide, 0.4, 6.9, 4, 0.3, "CONFIDENTIAL — Attorney Work Product", size=7, bold=True, color=HIGHLIGHT)
    add_text(slide, 10, 6.9, 3, 0.3, "2/2/26  |  Slide 2 of 3", size=7, color=GRAY, align=PP_ALIGN.RIGHT)

    # ============ SLIDE 3: Evidence Detail + Assessment ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK
    bar.line.fill.background()
    add_text(slide, 0.4, 0.1, 10, 0.35, f"Evidence & Assessment — {t['company']}", size=14, bold=True, color=WHITE)

    # Left column: element-by-element evidence
    y = 0.8
    for i, (element, (match, evidence)) in enumerate(zip(ELEMENTS, t["mapping"])):
        letter = element.split(")")[0] + ")"
        match_color = GREEN if match == "STRONG" else YELLOW_D if match == "LIKELY" else RED_D
        add_text(slide, 0.3, y, 0.4, 0.2, letter, size=8, bold=True, color=ACCENT)
        add_text(slide, 0.7, y, 0.8, 0.2, f"[{match}]", size=7, bold=True, color=match_color)
        add_text(slide, 1.6, y, 5.0, 0.35, evidence, size=7, color=DARK)
        y += 0.42

    # Right column: summary assessment
    add_text(slide, 7.2, 0.8, 3, 0.3, "INFRINGEMENT ASSESSMENT", size=10, bold=True, color=DARK)

    # Count matches
    strong = sum(1 for m, _ in t["mapping"] if m == "STRONG")
    likely = sum(1 for m, _ in t["mapping"] if m == "LIKELY")
    risk = sum(1 for m, _ in t["mapping"] if m == "RISK")

    stats = [
        f"STRONG matches: {strong}/7 elements",
        f"LIKELY matches: {likely}/7 elements",
        f"AT-RISK elements: {risk}/7 elements",
        f"Dependent claims met: {sum(1 for d in t['dep_match'] if d.startswith('YES'))}/4",
        "",
        f"Overall: {t['strength']} infringement likelihood",
    ]
    add_multiline(slide, 7.2, 1.2, 5.5, 1.5, stats, size=9, color=DARK, bold_first=False)

    # Next steps
    add_text(slide, 7.2, 3.0, 3, 0.3, "NEXT STEPS", size=10, bold=True, color=ACCENT)
    steps = [
        "1. Capture dated screenshots of all evidence URLs",
        "2. Record product demos (screen capture with timestamps)",
        "3. Review developer documentation / SDK / API docs",
        "4. Check target's own patent filings (shows field awareness)",
        "5. Identify technical experts for deposition if needed",
    ]
    add_multiline(slide, 7.2, 3.4, 5.8, 1.8, steps, size=8, color=DARK)

    # Sources (full URLs)
    add_text(slide, 7.2, 5.0, 3, 0.3, "EVIDENCE SOURCES", size=9, bold=True, color=ACCENT)
    add_multiline(slide, 7.2, 5.3, 5.8, 1.5, t["sources"], size=7, color=RGBColor(0, 80, 160))

    add_text(slide, 0.4, 6.9, 4, 0.3, "CONFIDENTIAL — Attorney Work Product", size=7, bold=True, color=HIGHLIGHT)
    add_text(slide, 10, 6.9, 3, 0.3, "2/2/26  |  Slide 3 of 3", size=7, color=GRAY, align=PP_ALIGN.RIGHT)

    # Save
    filename = f"EoU_Summary_{key}.pptx"
    filepath = os.path.join(OUTPUT_DIR, filename)
    prs.save(filepath)
    print(f"Created: {filepath}")


if __name__ == "__main__":
    for key, data in TARGETS.items():
        build_summary(key, data)
    print("\nDone — 4 summary decks created.")
