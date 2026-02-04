"""
Build Evidence of Use (EoU) PowerPoint files for WO/2025/193512A1
Following the Data Capture Fields format from USP 7062510 Experian exemplar.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))

PATENT_NUM = "WO/2025/193512A1"
APP_NUM = "PCT/US2025/XXXXXX"
PRIORITY_DATE = "2024 (estimated)"
TITLE = "SINGLE SHOT 3D MODELLING FROM 2D IMAGE"

CLAIM_PREAMBLE = "Claim 1: A method for generating a 3D model comprising:"

CLAIM_ELEMENTS = [
    "(a) obtaining a 2D image of an object;",
    "(b) classifying the object in the image;",
    "(c) segmenting the object from the image;",
    "(d) dimensionally sampling the segmented image of the object;",
    "(e) extracting texture information from the segmented image;",
    "(f) generating a 3D mesh model based at least on the dimensional sampling;",
    "(g) rendering the texture information onto the 3D mesh model to create a 3D representation of the object.",
]

# --- Target definitions ---

TARGETS = {
    "Google_Shopping": {
        "company": "Google (Alphabet Inc.)",
        "product": "Google Shopping 3D Product Visualizations",
        "elements": [
            {
                "evidence": (
                    "Google's 3D product viewer accepts as few as 1-3 product images as input. "
                    "Their Gen 3 pipeline (Veo-based) can generate 3D representations from a single image. "
                    "Google's research blog states: 'as few as three images capturing most object surfaces were sufficient.'"
                ),
                "source": "https://research.google/blog/bringing-3d-shoppable-products-online-with-generative-ai/",
                "commentary": (
                    "Google's system clearly obtains 2D images as the starting input. Their blog documents "
                    "accepting 1-3 images per product. This is the foundational input step of their pipeline."
                ),
            },
            {
                "evidence": (
                    "Google's system generalizes across diverse product categories including shoes, furniture, "
                    "apparel, electronics, and more. The Veo model was fine-tuned on 'millions of high-quality "
                    "3D synthetic assets' rendered by category. The system must classify input objects to select "
                    "appropriate generation parameters."
                ),
                "source": "https://research.google/blog/bringing-3d-shoppable-products-online-with-generative-ai/",
                "commentary": (
                    "While Google does not explicitly use the word 'classify,' their system handles diverse "
                    "product categories differently (shoes vs. furniture vs. electronics), which requires "
                    "object classification. The training pipeline is organized by product category. A formal "
                    "technical deposition or source code review could confirm the classification step."
                ),
            },
            {
                "evidence": (
                    "Google's blog confirms 'removing unwanted backgrounds' as a sub-problem in their pipeline. "
                    "The product is isolated from its background before 3D generation occurs. Gen 1 used "
                    "'object-centric images' — implying segmentation."
                ),
                "source": "https://research.google/blog/bringing-3d-shoppable-products-online-with-generative-ai/",
                "commentary": (
                    "Background removal is explicit segmentation — isolating the object from the image. "
                    "Google's own documentation confirms this step. This directly satisfies element (c)."
                ),
            },
            {
                "evidence": (
                    "Gen 1 used 'predicting 3D priors' and 'estimating camera positions from a sparse set of "
                    "object-centric images.' Gen 2 used Score Distillation Sampling where 'a 3D model is "
                    "rendered from a random camera view.' Gen 3 fine-tuned Veo on assets 'rendered from various "
                    "camera angles.' All generations extract dimensional/spatial information from segmented images."
                ),
                "source": "https://research.google/blog/bringing-3d-shoppable-products-online-with-generative-ai/",
                "commentary": (
                    "Dimensional sampling is the extraction of spatial/geometric information from the segmented "
                    "image. Google's pipeline extracts 3D priors, camera positions, and spatial relationships — "
                    "all forms of dimensional sampling. This is technically necessary for any 2D-to-3D conversion."
                ),
            },
            {
                "evidence": (
                    "Google's blog states Veo captures 'complex interactions between light, material, texture, "
                    "and geometry' and handles 'complex lighting and material interactions (i.e., shiny surfaces).' "
                    "Texture is explicitly extracted and preserved from the input images."
                ),
                "source": "https://research.google/blog/bringing-3d-shoppable-products-online-with-generative-ai/",
                "commentary": (
                    "Google explicitly mentions texture extraction as a core capability. The Veo model is "
                    "specifically designed to preserve texture fidelity. This directly maps to element (e)."
                ),
            },
            {
                "evidence": (
                    "Gen 1 used Neural Radiance Fields (NeRF) to generate 3D representations. Gen 2 optimized "
                    "'3D model parameters' via Score Distillation Sampling. Gen 3 generates 3D representations "
                    "via Veo. All generations produce a 3D geometric model from dimensional/spatial information."
                ),
                "source": "https://research.google/blog/bringing-3d-shoppable-products-online-with-generative-ai/",
                "commentary": (
                    "Google generates 3D representations (NeRF, mesh, or volumetric) from the spatial information "
                    "extracted from input images. Whether the intermediate representation is a NeRF, mesh, or "
                    "Gaussian splat, the output is functionally a 3D model. Google also supports merchant-uploaded "
                    "glTF/glb mesh files, confirming mesh-based output is part of their ecosystem."
                ),
            },
            {
                "evidence": (
                    "The final output is an interactive, photorealistic 360° product visualization on Google "
                    "Shopping. Products show realistic texture, lighting, and material properties. Users can "
                    "rotate and interact. Google states products with 3D imagery get '50% more user engagement.'"
                ),
                "source": "https://support.google.com/merchants/answer/13671720",
                "commentary": (
                    "The final rendered output is a textured 3D representation viewable from any angle. "
                    "This is the direct result of rendering texture onto the 3D model. The interactive "
                    "360° spin also satisfies dependent Claims 9, 10, and 11."
                ),
            },
        ],
    },
    "Meshy_AI": {
        "company": "Meshy Inc.",
        "product": "Meshy AI 3D Model Generator",
        "elements": [
            {
                "evidence": (
                    "Meshy accepts a single image upload as input. Their website states: 'Upload a clear image "
                    "or type a short text prompt' to generate a 3D model. Supported formats include standard "
                    "image files."
                ),
                "source": "https://www.meshy.ai/",
                "commentary": (
                    "Meshy's primary input mode is a single 2D image upload. This directly satisfies element (a). "
                    "The single-image workflow is the core product offering."
                ),
            },
            {
                "evidence": (
                    "Meshy generates 'front, side, and back view images' from the input before 3D conversion. "
                    "This multi-view generation step requires the system to understand what type of object it is "
                    "processing to generate coherent multi-view predictions. The system handles diverse object "
                    "categories (characters, products, environments)."
                ),
                "source": "https://www.meshy.ai/",
                "commentary": (
                    "The multi-view generation step inherently requires object classification — the system must "
                    "understand the object's category and structure to predict unseen views. Generating a coherent "
                    "back view of a shoe vs. a chair requires different learned priors per class."
                ),
            },
            {
                "evidence": (
                    "Meshy isolates the object from the background image to generate multi-view predictions "
                    "and ultimately the 3D mesh. The object must be segmented from the scene to produce a "
                    "clean 3D model without background artifacts."
                ),
                "source": "https://www.meshy.ai/",
                "commentary": (
                    "Object segmentation is a technically necessary step in Meshy's pipeline. A 3D model cannot "
                    "be generated from a cluttered scene without first isolating the target object. This is "
                    "standard practice in all image-to-3D systems."
                ),
            },
            {
                "evidence": (
                    "Meshy's 'Smart Remesh' feature allows adjustment of 'triangle or quad counts' ranging "
                    "'from 1k to 300k.' The system extracts geometric/dimensional information from the input "
                    "to produce mesh vertices and faces."
                ),
                "source": "https://www.meshy.ai/",
                "commentary": (
                    "The mesh generation with adjustable polygon counts demonstrates that dimensional/spatial "
                    "sampling occurs — the system samples geometric information from the segmented image to "
                    "construct 3D vertices. The adjustable resolution directly shows dimensional sampling at "
                    "different granularities."
                ),
            },
            {
                "evidence": (
                    "Meshy generates full PBR (Physically Based Rendering) texture maps: 'Diffuse, Roughness, "
                    "Metallic, and Normal maps.' Multiple texture styles available including 'realistic and "
                    "cartoon to hand painted and fantasy.' Textures are extracted/generated from the input image."
                ),
                "source": "https://www.meshy.ai/",
                "commentary": (
                    "Meshy explicitly extracts and generates texture information from the input image. The PBR "
                    "texture pipeline (Diffuse, Roughness, Metallic, Normal) is a comprehensive texture "
                    "extraction system. This directly and unambiguously satisfies element (e)."
                ),
            },
            {
                "evidence": (
                    "Meshy generates explicit 3D meshes exportable in 'FBX, GLB, OBJ, STL, 3MF, USDZ, BLEND "
                    "formats.' The Smart Remesh feature produces meshes with configurable triangle/quad counts. "
                    "This is an explicit mesh generation step."
                ),
                "source": "https://www.meshy.ai/",
                "commentary": (
                    "Meshy's output is unambiguously a 3D mesh model. The multiple mesh export formats (OBJ, "
                    "FBX, GLB, STL) and adjustable polygon counts confirm explicit mesh generation. This is "
                    "the strongest possible match for element (f) — there is no ambiguity."
                ),
            },
            {
                "evidence": (
                    "The final output is a textured 3D mesh viewable in Meshy's browser-based 3D viewer. "
                    "PBR textures are rendered onto the mesh geometry. Users can rotate, zoom, and inspect "
                    "the model from any angle before downloading."
                ),
                "source": "https://www.meshy.ai/",
                "commentary": (
                    "Texture rendering onto the mesh is the final pipeline step. The browser-based 3D viewer "
                    "demonstrates the textured 3D representation. This also satisfies Claims 9 (multiple views), "
                    "10 (user input rotation), and 11 (view capture via export)."
                ),
            },
        ],
    },
    "Apple_SHARP": {
        "company": "Apple Inc.",
        "product": "SHARP (Sharp Monocular View Synthesis) + Spatial Scenes (iOS 26)",
        "elements": [
            {
                "evidence": (
                    "SHARP takes a single photograph as input. Apple's documentation: 'Given a single photograph, "
                    "it regresses the parameters of a 3D Gaussian representation.' The CLI command is simply: "
                    "'sharp predict -i image_path.png'"
                ),
                "source": "https://github.com/apple/ml-sharp",
                "commentary": (
                    "SHARP's entire design premise is single-image input. The GitHub repo, research paper, and "
                    "Apple's ML research page all confirm this. This directly satisfies element (a)."
                ),
            },
            {
                "evidence": (
                    "SHARP demonstrates 'robust zero-shot generalization across datasets' — it works on diverse "
                    "object and scene categories without per-category fine-tuning. The model was trained on data "
                    "spanning multiple object and scene types."
                ),
                "source": "https://machinelearning.apple.com/research/sharp-monocular-view",
                "commentary": (
                    "While SHARP uses zero-shot generalization rather than explicit classification, the neural "
                    "network inherently learns object category representations during training. The model must "
                    "internally differentiate object types to generate appropriate 3D structures. Source code "
                    "review could confirm whether explicit classification occurs in the architecture."
                ),
            },
            {
                "evidence": (
                    "SHARP uses Apple's Depth Pro for depth estimation, which requires isolating scene elements "
                    "and understanding object boundaries. The model predicts per-pixel 3D Gaussian parameters, "
                    "which requires understanding object vs. background separation."
                ),
                "source": "https://github.com/apple/ml-sharp",
                "commentary": (
                    "Depth estimation inherently involves segmentation — distinguishing foreground objects from "
                    "background to assign correct depth values. Depth Pro's per-pixel predictions require "
                    "implicit object segmentation. This satisfies element (c), though explicit segmentation "
                    "masks may or may not be generated as an intermediate step."
                ),
            },
            {
                "evidence": (
                    "SHARP predicts 'position and appearance of 3D Gaussians' — each Gaussian has a 3D position, "
                    "scale, and orientation. The output is metric-scale with 'absolute scale, supporting metric "
                    "camera movements.' Depth Pro provides dimensional depth information."
                ),
                "source": "https://machinelearning.apple.com/research/sharp-monocular-view",
                "commentary": (
                    "The prediction of 3D Gaussian positions with metric scale is a form of dimensional sampling — "
                    "the system samples the spatial dimensions of the object from the 2D image. Each Gaussian's "
                    "position, scale, and orientation encode dimensional information."
                ),
            },
            {
                "evidence": (
                    "SHARP predicts 'appearance' parameters for each 3D Gaussian, encoding color and texture "
                    "information from the input image. The output achieves 'photorealistic' rendering quality "
                    "with 'sharp details and fine structures.'"
                ),
                "source": "https://apple.github.io/ml-sharp/",
                "commentary": (
                    "The 'appearance' parameters of each Gaussian encode texture/color information extracted "
                    "from the input image. This is functionally equivalent to texture extraction — the system "
                    "captures visual appearance data from the segmented image."
                ),
            },
            {
                "evidence": (
                    "SHARP generates a 3D Gaussian Splatting representation (.ply format). The output is a "
                    "structured 3D representation with explicit positions, scales, rotations, and spherical "
                    "harmonic coefficients. Compatible with standard 3DGS renderers."
                ),
                "source": "https://huggingface.co/apple/Sharp",
                "commentary": (
                    "POTENTIAL DESIGN-AROUND: SHARP outputs 3D Gaussian splats, not a traditional polygon mesh. "
                    "The claim language says '3D mesh model.' Whether a Gaussian splat representation reads on "
                    "'mesh model' is a claim construction question. Argument FOR infringement: the functional "
                    "result is identical (a 3D model of the object). Argument AGAINST: Gaussians are not "
                    "vertices/faces/edges. However, Gaussian splats can be converted to meshes, and Apple's "
                    "ecosystem (Object Capture, RealityKit) works with mesh formats."
                ),
            },
            {
                "evidence": (
                    "SHARP renders photorealistic novel views at '100+ frames per second on a standard GPU.' "
                    "The textured 3D representation can be viewed from nearby viewpoints with 'sharp details "
                    "and fine structures.' Already integrated into Splat Studio app for Vision Pro."
                ),
                "source": "https://apple.github.io/ml-sharp/",
                "commentary": (
                    "The rendering of textured 3D content from multiple viewpoints satisfies element (g) and "
                    "dependent Claims 9 and 11. The Vision Pro integration via Splat Studio demonstrates "
                    "commercial deployment. Apple also commercialized related tech as 'Spatial Scenes' in iOS 26, "
                    "strengthening the commercial infringement argument."
                ),
            },
        ],
    },
    "Amazon": {
        "company": "Amazon.com, Inc.",
        "product": "Amazon 3D/AR Product Experience (Seller App 3D Scanning + View in 3D + AWS 3D Pipeline)",
        "elements": [
            {
                "evidence": (
                    "Amazon Seller App 'Create 3D Models' tool captures product images via iOS device camera. "
                    "AWS Human Mesh Recovery pipeline processes '2D images or videos' as input. Amazon's 3D "
                    "listing requirement: '2-10 reference photos' submitted with product dimensions. "
                    "Third-party providers (Nextech3D.ai) generate 3D models from as few as 1-2 product photos."
                ),
                "source": "https://sell.amazon.com/tools/3d-ar",
                "commentary": (
                    "Amazon's ecosystem obtains 2D images at multiple points: seller scanning, reference photo "
                    "upload, and third-party provider input. The entire 3D pipeline begins with 2D image capture. "
                    "Whether Amazon's own scanning tool or their certified provider Nextech3D.ai is used, "
                    "the input is always 2D images."
                ),
            },
            {
                "evidence": (
                    "Amazon's 3D features are category-dependent: 'Eligibility depends on product category, "
                    "physical properties, and other factors.' Amazon categorizes products (shoes, eyewear, "
                    "furniture, electronics) and applies different 3D experiences per category: Virtual Try-On "
                    "for shoes/eyewear, View in Your Room for furniture, View in 3D for general products."
                ),
                "source": "https://sell.amazon.com/tools/3d-ar",
                "commentary": (
                    "Amazon explicitly classifies products by category to determine which 3D experience to "
                    "apply. This product classification step directly maps to element (b). The system must "
                    "understand the object type to generate an appropriate 3D model and select the correct "
                    "viewing experience."
                ),
            },
            {
                "evidence": (
                    "Amazon's 3D scanning isolates the product from its environment. AWS 3D Gaussian Splatting "
                    "blog describes using 'image segmentation and semantic labelling' via Amazon SageMaker "
                    "and Amazon Bedrock. The product object must be segmented from the scene."
                ),
                "source": "https://aws.amazon.com/blogs/spatial/3d-gaussian-splatting-performant-3d-scene-reconstruction-at-scale/",
                "commentary": (
                    "AWS explicitly mentions 'image segmentation' as part of their 3D reconstruction pipeline. "
                    "This is documented in their own technical blog. The Seller App scanning also requires "
                    "product isolation from the scanning environment."
                ),
            },
            {
                "evidence": (
                    "Amazon requires 'accurate product dimensions' submitted with 3D models. AWS Human Mesh "
                    "Recovery uses ScoreHMR with 'diffusion models to capture and reconstruct human body "
                    "parameters from input images.' The 3D scanning tool extracts spatial/dimensional data "
                    "from captured images over a 5-10 minute scanning session."
                ),
                "source": "https://aws.amazon.com/blogs/spatial/from-2d-to-3d-building-a-scalable-human-mesh-recovery-pipeline-with-amazon-sagemaker-ai/",
                "commentary": (
                    "Dimensional sampling occurs during the scanning process (extracting 3D geometry from "
                    "multiple 2D captures) and in the AWS pipeline (extracting body parameters from images). "
                    "The requirement for 'accurate product dimensions' confirms dimensional data extraction."
                ),
            },
            {
                "evidence": (
                    "Amazon's 3D product models display photorealistic textures — product colors, labels, "
                    "branding, and material properties are preserved. The scanning process captures visual "
                    "appearance data. AWS 3D pipeline handles texture as part of reconstruction."
                ),
                "source": "https://sell.amazon.com/tools/3d-ar",
                "commentary": (
                    "Texture extraction is inherent in any 3D product visualization system that produces "
                    "photorealistic results. Amazon's 3D product views clearly show extracted texture "
                    "information from the original product images."
                ),
            },
            {
                "evidence": (
                    "Amazon requires 3D models in GLB or GLTF format — both are mesh-based 3D formats "
                    "(vertices, faces, materials). AWS blog describes 'mesh representations' and 'human mesh "
                    "recovery.' The Seller App produces 3D models reviewed by '3D artists' before publication."
                ),
                "source": "https://sell.amazon.com/tools/3d-ar",
                "commentary": (
                    "Amazon's 3D ecosystem is explicitly mesh-based. GLB/GLTF are mesh formats by definition. "
                    "This is the strongest match for element (f) — Amazon requires and generates 3D mesh models. "
                    "No ambiguity about the representation format."
                ),
            },
            {
                "evidence": (
                    "Amazon's 'View in 3D' renders textured 3D models allowing customers to 'rotate and zoom in "
                    "on a product in 3D, offering a detailed view from all angles.' 'View in Your Room' renders "
                    "the textured model in AR. Amazon reports '2X improvement in purchase conversion' and "
                    "'20% lower return rate' from 3D/AR features."
                ),
                "source": "https://sell.amazon.com/tools/3d-ar",
                "commentary": (
                    "The final rendered output is a textured 3D representation viewable from any angle. "
                    "This satisfies element (g) and dependent Claims 9 (multiple views), 10 (user rotation), "
                    "and 11 (2D captures from each view). Amazon's massive deployment scale (millions of "
                    "product listings) makes this a high-value target."
                ),
            },
        ],
    },
}


def add_footer(slide, slide_num, date="2/2/26"):
    """Add date, CONFIDENTIAL, and slide number footer."""
    # Date
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(6.8), Inches(1.5), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = date
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(128, 128, 128)

    # CONFIDENTIAL
    txBox = slide.shapes.add_textbox(Inches(4), Inches(6.8), Inches(2), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "CONFIDENTIAL"
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = RGBColor(200, 0, 0)
    p.alignment = PP_ALIGN.CENTER

    # Slide number
    txBox = slide.shapes.add_textbox(Inches(8.5), Inches(6.8), Inches(1), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(slide_num)
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(128, 128, 128)
    p.alignment = PP_ALIGN.RIGHT


def add_textbox(slide, left, top, width, height, text, font_size=11, bold=False, color=None, alignment=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    if bold:
        p.font.bold = True
    if color:
        p.font.color.rgb = color
    if alignment:
        p.alignment = alignment
    return txBox


def build_pptx(target_key, target_data):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    company = target_data["company"]
    product = target_data["product"]
    elements = target_data["elements"]

    # ---- SLIDE 1: Title ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_textbox(slide, 0.5, 1.0, 9, 1.2,
                f"Mapping of {PATENT_NUM} on\n{product}",
                font_size=22, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 2.8, 9, 0.5,
                f"{TITLE}",
                font_size=14, alignment=PP_ALIGN.CENTER,
                color=RGBColor(80, 80, 80))
    add_textbox(slide, 0.5, 3.8, 4, 0.8,
                f"Priority: {PRIORITY_DATE}",
                font_size=11, color=RGBColor(100, 100, 100))
    add_textbox(slide, 0.5, 5.0, 4, 0.4,
                "Field 1: Patent Number", font_size=9,
                color=RGBColor(0, 0, 200))
    add_textbox(slide, 0.5, 5.4, 4, 0.4,
                f"Field 9-10: Company: {company} | Product: {product}",
                font_size=9, color=RGBColor(0, 0, 200))
    add_footer(slide, 1)

    # ---- SLIDE 2: Claim Preamble ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Claim preamble
    add_textbox(slide, 0.3, 0.3, 5.5, 0.8,
                CLAIM_PREAMBLE,
                font_size=12, bold=True)
    # Field labels
    add_textbox(slide, 6.2, 0.3, 3.5, 0.5,
                "Field 2: Claim Number\nFocus on Claim 1 (independent)",
                font_size=9, color=RGBColor(0, 0, 200))
    add_textbox(slide, 6.2, 1.0, 3.5, 0.4,
                "Field 3: Claim Preamble",
                font_size=9, color=RGBColor(0, 0, 200))

    # Evidence area label
    add_textbox(slide, 0.3, 1.5, 5.5, 0.4,
                "Field 5: Evidence 1 — Screenshots or copy/pastes from official sources",
                font_size=9, color=RGBColor(0, 0, 200))
    add_textbox(slide, 0.3, 1.9, 5.5, 0.4,
                "Field 6: Evidence 2 — Additional evidence with highlighted infringement",
                font_size=9, color=RGBColor(0, 0, 200))
    add_textbox(slide, 0.3, 2.3, 5.5, 0.4,
                "Field 7: Source — URLs tied to evidence",
                font_size=9, color=RGBColor(0, 0, 200))

    # Commentary explanation
    add_textbox(slide, 0.3, 3.0, 9, 2.5,
                "Field 8: Student Commentary\n\n"
                "For each claim element, provide:\n"
                "(1) Why you believe there is infringement\n"
                "(2) How you know, either:\n"
                "    (a) It's the only way it could be done technically, or\n"
                "    (b) We'd need XYZ (source code, etc.) to prove it\n\n"
                f"Target Company: {company}\n"
                f"Target Product: {product}",
                font_size=11)
    add_footer(slide, 2)

    # ---- SLIDES 3+: One per claim element ----
    for i, (element_text, evidence_data) in enumerate(zip(CLAIM_ELEMENTS, elements)):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Field 4: Claim Limitation
        add_textbox(slide, 0.3, 0.2, 9.2, 0.3,
                    "Field 4: Claim Limitation (Element)",
                    font_size=8, color=RGBColor(0, 0, 200))
        add_textbox(slide, 0.3, 0.5, 9.2, 0.7,
                    element_text,
                    font_size=12, bold=True)

        # Evidence box
        add_textbox(slide, 0.3, 1.4, 9.2, 0.3,
                    "Field 5/6: Evidence",
                    font_size=8, color=RGBColor(0, 0, 200))

        # Evidence content
        ev_box = add_textbox(slide, 0.3, 1.7, 9.2, 1.8,
                             evidence_data["evidence"],
                             font_size=10)
        # Add border-like background
        ev_box.fill.solid()
        ev_box.fill.fore_color.rgb = RGBColor(245, 245, 245)

        # Source URL
        add_textbox(slide, 0.3, 3.7, 9.2, 0.3,
                    "Field 7: Source",
                    font_size=8, color=RGBColor(0, 0, 200))
        add_textbox(slide, 0.3, 4.0, 9.2, 0.3,
                    evidence_data["source"],
                    font_size=9, color=RGBColor(0, 100, 200))

        # Commentary
        add_textbox(slide, 0.3, 4.5, 9.2, 0.3,
                    "Field 8: Commentary",
                    font_size=8, color=RGBColor(0, 0, 200))
        comm_box = add_textbox(slide, 0.3, 4.8, 9.2, 1.8,
                               evidence_data["commentary"],
                               font_size=10)
        comm_box.fill.solid()
        comm_box.fill.fore_color.rgb = RGBColor(255, 255, 235)

        add_footer(slide, i + 3)

    # Save
    filename = f"EoU_{PATENT_NUM.replace('/', '_')}_{target_key}.pptx"
    filepath = os.path.join(OUTPUT_DIR, filename)
    prs.save(filepath)
    print(f"Created: {filepath}")
    return filepath


if __name__ == "__main__":
    for key, data in TARGETS.items():
        build_pptx(key, data)
    print("\nDone. All 4 EoU files created.")
