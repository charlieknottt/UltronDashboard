"""
Concise EoU PowerPoints for WO/2025/193512A1.
No template labels. No filler. Only strong evidence per target.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))
PATENT = "WO/2025/193512A1"


def add_text(slide, left, top, w, h, text, size=11, bold=False, color=RGBColor(0,0,0), align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_para(tf, text, size=9, bold=False, color=RGBColor(0,0,0), space_after=Pt(2)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.space_after = space_after
    return p


def screenshot_box(slide, left, top, w, h):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 240, 243)
    shape.line.color.rgb = RGBColor(190, 190, 200)
    shape.line.dash_style = 2
    shape.line.width = Pt(0.75)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "[Screenshot]"
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(160, 160, 170)
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER


def footer(slide, num, total):
    add_text(slide, 0.3, 7.05, 3, 0.2, "CONFIDENTIAL", size=7, bold=True, color=RGBColor(180,0,0))
    add_text(slide, 8.2, 7.05, 1.5, 0.2, f"{num}/{total}", size=7, color=RGBColor(150,150,150), align=PP_ALIGN.RIGHT)


def build(filename, company, product, slides_data):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    total = len(slides_data) + 1  # +1 for title

    # Title slide
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(20, 20, 45); bar.line.fill.background()

    add_text(sl, 0.5, 1.5, 9, 0.6, f"Evidence of Use: {PATENT}", size=11, color=RGBColor(100,100,120))
    add_text(sl, 0.5, 2.1, 9, 0.8, f"{company}", size=26, bold=True, color=RGBColor(20,20,45))
    add_text(sl, 0.5, 2.9, 9, 0.5, product, size=13, color=RGBColor(60,60,80))
    add_text(sl, 0.5, 4.0, 9, 0.3, "SINGLE SHOT 3D MODELLING FROM 2D IMAGE", size=10, color=RGBColor(130,130,150))
    add_text(sl, 0.5, 4.4, 9, 0.3, "Carnegie Mellon University  |  Savvides et al.", size=10, color=RGBColor(130,130,150))
    footer(sl, 1, total)

    # Element slides
    for i, sd in enumerate(slides_data):
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.08))
        bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(20, 20, 45); bar.line.fill.background()

        # Claim element
        add_text(sl, 0.3, 0.25, 9.4, 0.5, sd["claim"], size=12, bold=True, color=RGBColor(20,20,45))

        # Evidence (left side)
        screenshot_box(sl, 0.3, 0.9, 4.4, 2.6)

        # Evidence text (right side)
        box = add_text(sl, 4.9, 0.9, 4.8, 2.6, "", size=1)
        tf = box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = ""
        for line in sd["evidence"]:
            add_para(tf, line, size=9, color=RGBColor(30,30,30), space_after=Pt(4))

        # Source URLs
        add_text(sl, 0.3, 3.7, 1, 0.2, "Source:", size=8, bold=True, color=RGBColor(80,80,100))
        y = 3.7
        for src in sd["sources"]:
            add_text(sl, 1.2, y, 8.5, 0.2, src, size=8, color=RGBColor(0, 70, 170))
            y += 0.2

        # Second screenshot area
        screenshot_box(sl, 0.3, y + 0.2, 4.4, 1.6)

        # Commentary
        cy = y + 0.2
        cbox = add_text(sl, 4.9, cy, 4.8, 1.6, "", size=1)
        ctf = cbox.text_frame
        ctf.word_wrap = True
        ctf.paragraphs[0].text = ""
        add_para(ctf, "Commentary", size=8, bold=True, color=RGBColor(80,80,100))
        add_para(ctf, sd["commentary"], size=9, color=RGBColor(30,30,30))

        footer(sl, i + 2, total)

    path = os.path.join(OUTPUT_DIR, filename)
    prs.save(path)
    print(f"{filename} — {total} slides")


# =====================================================================
# GOOGLE — focus on their blog documenting the exact pipeline
# =====================================================================
build("EoU_Google_Shopping.pptx", "Google (Alphabet)", "Google Shopping — AI-Generated 3D Product Viewer", [
    {
        "claim": "(a) obtaining a 2D image  +  (c) segmenting the object",
        "evidence": [
            "Google's Gen 3 pipeline accepts as few as 1 product image.",
            "Their research blog confirms 'removing unwanted backgrounds' as an explicit pipeline step.",
            "Product is isolated from its scene before 3D generation begins.",
        ],
        "sources": ["https://research.google/blog/bringing-3d-shoppable-products-online-with-generative-ai/"],
        "commentary": "Google's own blog documents both the image input requirement and background removal (segmentation). Direct public admission of both elements.",
    },
    {
        "claim": "(d) dimensionally sampling  +  (f) generating a 3D model",
        "evidence": [
            "Gen 1: NeRF with '3D priors prediction' and 'camera position estimation from sparse images.'",
            "Gen 2: Score Distillation Sampling — '3D model rendered from random camera view, diffusion model generates target from same view, score informs 3D model parameter optimization.'",
            "Gen 3: Veo fine-tuned on 'millions of 3D synthetic assets rendered from various camera angles.'",
        ],
        "sources": ["https://research.google/blog/bringing-3d-shoppable-products-online-with-generative-ai/"],
        "commentary": "All three generations extract spatial dimensions from 2D images and produce 3D representations. The blog contains pipeline diagrams showing each step. Three generations of evidence.",
    },
    {
        "claim": "(e) extracting texture  +  (g) rendering texture onto 3D representation",
        "evidence": [
            "Blog: Veo captures 'complex interactions between light, material, texture, and geometry.'",
            "Handles 'complex lighting and material interactions (i.e., shiny surfaces).'",
            "Output: interactive photorealistic 360° spins on Google Shopping — users rotate product freely.",
            "Google reports 50% more engagement on products with 3D imagery.",
        ],
        "sources": [
            "https://research.google/blog/bringing-3d-shoppable-products-online-with-generative-ai/",
            "https://support.google.com/merchants/answer/13671720",
        ],
        "commentary": "Google explicitly names 'texture' in their technical description. The rendered 360° output also satisfies Claims 9 (multiple views), 10 (user rotation), and 11 (frame capture).",
    },
])


# =====================================================================
# MESHY — cleanest match, explicit mesh + PBR textures
# =====================================================================
build("EoU_Meshy_AI.pptx", "Meshy Inc.", "Meshy AI — Image-to-3D Model Generator", [
    {
        "claim": "(a) obtaining a 2D image  +  (c) segmenting the object",
        "evidence": [
            "Core feature: single image upload → 3D model.",
            "System generates front/side/back views from the single image — object must be isolated from background to produce clean multi-view predictions.",
            "Output 3D models contain only the target object, never background elements.",
        ],
        "sources": ["https://www.meshy.ai/", "https://www.meshy.ai/features/image-to-3d"],
        "commentary": "The output itself proves segmentation: every generated 3D model is a clean isolated object with no background. Input is unambiguously a single 2D image.",
    },
    {
        "claim": "(d) dimensionally sampling  +  (f) generating a 3D mesh model",
        "evidence": [
            "Smart Remesh: adjustable polygon count from 1k to 300k triangles/quads.",
            "Exports in 7 mesh formats: FBX, GLB, OBJ, STL, 3MF, USDZ, BLEND.",
            "These are all standard polygon mesh formats (vertices, faces, edges) — no ambiguity.",
        ],
        "sources": ["https://www.meshy.ai/"],
        "commentary": "Strongest possible match for element (f). OBJ, FBX, GLB, STL are mesh formats by definition. Adjustable polygon counts demonstrate dimensional sampling at configurable granularity.",
    },
    {
        "claim": "(e) extracting texture  +  (g) rendering texture onto mesh",
        "evidence": [
            "Full PBR texture pipeline: Diffuse, Roughness, Metallic, and Normal maps.",
            "Multiple texture styles: 'realistic, cartoon, hand painted, fantasy.'",
            "Browser-based 3D viewer renders textures onto mesh — users rotate and inspect from any angle.",
            "Separate AI Texturing feature re-textures existing models.",
        ],
        "sources": ["https://www.meshy.ai/", "https://www.meshy.ai/features/ai-texturing"],
        "commentary": "PBR texture maps (Diffuse, Roughness, Metallic, Normal) are comprehensive texture extraction. No ambiguity on elements (e) or (g). Also satisfies Claims 9, 10, 11 via interactive viewer + export.",
    },
    {
        "claim": "Claim 5: 3D mesh generated by trained neural network",
        "evidence": [
            "Meshy's entire pipeline is AI/neural network-based.",
            "Image-to-3D uses trained models to predict 3D geometry from 2D input.",
            "1M+ users, 50+ concurrent model generation tasks — commercial scale.",
        ],
        "sources": ["https://www.meshy.ai/"],
        "commentary": "Claim 5 requires the mesh be generated by a 'trained neural network' — Meshy's core technology. Their scale (1M+ users) establishes significant commercial value for licensing.",
    },
])


# =====================================================================
# APPLE — SHARP is open-source, fully documented, but Gaussian splat risk
# =====================================================================
build("EoU_Apple_SHARP.pptx", "Apple Inc.", "SHARP + Spatial Scenes (iOS 26) + Vision Pro", [
    {
        "claim": "(a) obtaining a 2D image",
        "evidence": [
            "SHARP: 'Given a single photograph, regresses parameters of a 3D Gaussian representation.'",
            "CLI: 'sharp predict -i image_path.png' — single image, single command.",
            "Published on GitHub (apple/ml-sharp), Hugging Face, and Apple ML Research.",
        ],
        "sources": [
            "https://github.com/apple/ml-sharp",
            "https://machinelearning.apple.com/research/sharp-monocular-view",
        ],
        "commentary": "Entire system designed for single-image input. Fully documented across GitHub, Hugging Face, and Apple's research page. No ambiguity.",
    },
    {
        "claim": "(d) dimensionally sampling  +  (e) extracting texture",
        "evidence": [
            "Predicts per-pixel 3D Gaussian positions with metric scale ('absolute scale, supporting metric camera movements').",
            "Each Gaussian has position, scale, rotation (dimensional) + spherical harmonic appearance coefficients (texture).",
            "Uses Apple Depth Pro for monocular depth estimation — extracts depth dimensions per pixel.",
            "Output achieves 'photorealistic' quality with 'sharp details and fine structures.'",
        ],
        "sources": [
            "https://machinelearning.apple.com/research/sharp-monocular-view",
            "https://huggingface.co/papers/2512.10685",
        ],
        "commentary": "Metric-scale 3D position prediction IS dimensional sampling. Spherical harmonic coefficients encode texture/color. Both elements clearly satisfied by the documented architecture.",
    },
    {
        "claim": "(f) generating a 3D model  —  CLAIM CONSTRUCTION ISSUE",
        "evidence": [
            "SHARP outputs 3D Gaussian Splats (.ply format), NOT a traditional polygon mesh.",
            "Patent language: '3D mesh model.'",
            "Gaussian splats are point-based representations (no vertices/faces/edges).",
            "However: functionally equivalent (3D representation of object), convertible to mesh, and Apple's ecosystem (RealityKit, Object Capture) works with mesh formats.",
        ],
        "sources": ["https://github.com/apple/ml-sharp", "https://huggingface.co/apple/Sharp"],
        "commentary": "Whether '3D Gaussian splat' reads on '3D mesh model' is the key claim construction question. Functional equivalence argument is strong but not certain. Patent counsel should evaluate. This is the weakest element in the Apple case.",
    },
    {
        "claim": "(g) rendering texture onto 3D representation  +  Commercial deployment",
        "evidence": [
            "Renders photorealistic novel views at 100+ FPS on standard GPU.",
            "Integrated into Splat Studio app for Apple Vision Pro.",
            "Apple commercialized related technology as 'Spatial Scenes' in iOS 26 (2025).",
            "Open-source code on GitHub — anyone can verify the pipeline.",
        ],
        "sources": [
            "https://apple.github.io/ml-sharp/",
            "https://www.uploadvr.com/apple-sharp-open-source-on-device-gaussian-splatting/",
        ],
        "commentary": "Spatial Scenes in iOS 26 moves this from research to commercial deployment across every iPhone. Satisfies Claims 9 and 10 (novel view rendering with user control). Open-source repo enables full technical verification.",
    },
])


# =====================================================================
# AMAZON — mandatory 3D since Dec 2023, GLB/GLTF = unambiguous mesh
# =====================================================================
build("EoU_Amazon.pptx", "Amazon.com, Inc.", "Amazon 3D/AR — Seller App + View in 3D + AWS Pipeline", [
    {
        "claim": "(a) obtaining a 2D image  +  (b) classifying the object",
        "evidence": [
            "Seller App 'Create 3D Models': captures product images via iOS camera.",
            "Requires 2-10 reference photos + accurate product dimensions for 3D submission.",
            "Explicit product classification: 'Eligibility depends on product category, physical properties.'",
            "Different 3D experiences per category — Virtual Try-On (shoes), View in Your Room (furniture), View in 3D (general).",
        ],
        "sources": [
            "https://sell.amazon.com/tools/3d-ar",
        ],
        "commentary": "Amazon's pipeline starts with 2D images and classifies products by category to determine which 3D experience applies. Category-dependent eligibility is documented on their seller page.",
    },
    {
        "claim": "(c) segmenting  +  (d) dimensionally sampling",
        "evidence": [
            "AWS Spatial Computing blog: 'image segmentation and semantic labelling' via SageMaker and Bedrock.",
            "AWS Human Mesh Recovery: ScoreHMR uses 'diffusion models to capture and reconstruct body parameters from input images.'",
            "Seller App scanning extracts 3D spatial data over 5-10 minute session.",
            "'Accurate product dimensions' required with every 3D model submission.",
        ],
        "sources": [
            "https://aws.amazon.com/blogs/spatial/3d-gaussian-splatting-performant-3d-scene-reconstruction-at-scale/",
            "https://aws.amazon.com/blogs/spatial/from-2d-to-3d-building-a-scalable-human-mesh-recovery-pipeline-with-amazon-sagemaker-ai/",
        ],
        "commentary": "AWS explicitly documents 'image segmentation' in their own 3D pipeline blog — direct admission. Dimensional requirements for submissions confirm spatial data extraction.",
    },
    {
        "claim": "(f) generating a 3D mesh model",
        "evidence": [
            "All 3D models must be GLB or GLTF format — mesh formats by specification (vertices, faces, normals, materials).",
            "Amazon mandated 3D models over 360° images starting December 14, 2023.",
            "3D artist review of every scanned model before publication.",
        ],
        "sources": [
            "https://sell.amazon.com/tools/3d-ar",
        ],
        "commentary": "GLB/GLTF ARE mesh formats. No design-around argument possible. Amazon mandates mesh output for all 3D product content across the platform.",
    },
    {
        "claim": "(g) rendering texture onto 3D mesh  +  Claims 9/10/11",
        "evidence": [
            "View in 3D: 'rotate and zoom in on a product in 3D from all angles.'",
            "View in Your Room: AR rendering of textured 3D product in user's space.",
            "Virtual Try-On: renders product directly on user (shoes, eyewear).",
            "Impact: 2X purchase conversion, 20% lower returns, 8X increase in AR usage (2018-2022).",
        ],
        "sources": [
            "https://sell.amazon.com/tools/3d-ar",
        ],
        "commentary": "Three distinct rendering experiences (3D viewer, AR room, virtual try-on) all satisfy element (g) and dependent Claims 9-11. Amazon's own conversion data quantifies the commercial value of the infringing feature.",
    },
])

print("\nAll 4 final EoU decks created.")
