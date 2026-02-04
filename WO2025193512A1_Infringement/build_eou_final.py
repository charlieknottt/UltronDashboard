"""
Build EoU PowerPoints matching the Experian exemplar format exactly.
Each claim element gets its own slide with:
  - Claim limitation text (top)
  - Evidence area (screenshot placeholder + description)
  - Source URL (for screenshotting)
  - Commentary (why it infringes + what proof is needed)

Variable length per target — only elements with real evidence get slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))

PATENT = "WO/2025/193512A1"


def add_text(slide, left, top, w, h, text, size=11, bold=False, color=RGBColor(0,0,0), align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_screenshot_placeholder(slide, left, top, w, h, label="[INSERT SCREENSHOT HERE]"):
    """Gray dashed box where screenshot should go."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(235, 235, 240)
    shape.line.color.rgb = RGBColor(180, 180, 190)
    shape.line.dash_style = 2  # dash
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(140, 140, 150)
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER


def add_footer(slide, slide_num, date="2/2/26"):
    add_text(slide, 0.3, 7.0, 1.2, 0.25, date, size=7, color=RGBColor(150,150,150))
    add_text(slide, 4.0, 7.0, 2, 0.25, "CONFIDENTIAL", size=7, bold=True, color=RGBColor(200,0,0), align=PP_ALIGN.CENTER)
    add_text(slide, 8.5, 7.0, 1.2, 0.25, str(slide_num), size=7, color=RGBColor(150,150,150), align=PP_ALIGN.RIGHT)


def build_eou(filename, company, product, elements):
    """
    elements: list of dicts, each with:
      - claim_text: the claim limitation language
      - evidence_desc: text description of evidence
      - screenshot_label: what to screenshot (instructions)
      - sources: list of URLs to screenshot from
      - commentary: why it infringes + proof needed
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_num = 1

    # ---- SLIDE 1: Title ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, 0.5, 0.8, 9, 0.8,
             f"Mapping of {PATENT} on\n{product}",
             size=22, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 0.5, 2.0, 9, 0.4,
             "SINGLE SHOT 3D MODELLING FROM 2D IMAGE",
             size=12, color=RGBColor(100,100,100), align=PP_ALIGN.CENTER)
    add_text(slide, 0.5, 3.0, 9, 0.4,
             f"Target Company: {company}",
             size=11, align=PP_ALIGN.CENTER)
    add_text(slide, 0.5, 3.4, 9, 0.4,
             f"Target Product: {product}",
             size=11, align=PP_ALIGN.CENTER)
    add_text(slide, 0.5, 4.2, 4, 0.3,
             "Carnegie Mellon University", size=10, color=RGBColor(100,100,100))
    add_text(slide, 0.5, 4.5, 4, 0.3,
             "Inventor: Marios Savvides et al.", size=10, color=RGBColor(100,100,100))
    add_text(slide, 0.5, 5.5, 4, 0.3,
             "Field 1: Patent Number", size=9, color=RGBColor(0,0,200))
    add_text(slide, 0.5, 5.8, 9, 0.3,
             f"Field 9/10: Company: {company}  |  Product: {product}",
             size=9, color=RGBColor(0,0,200))
    add_footer(slide, slide_num)
    slide_num += 1

    # ---- SLIDE 2: Claim Preamble ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, 0.3, 0.3, 6, 0.5,
             "Claim 1: A method for generating a 3D model comprising:",
             size=13, bold=True)
    add_text(slide, 6.5, 0.3, 3.2, 0.5,
             "Field 2: Claim Number — Claim 1",
             size=9, color=RGBColor(0,0,200))
    add_text(slide, 6.5, 0.7, 3.2, 0.3,
             "Field 3: Claim Preamble",
             size=9, color=RGBColor(0,0,200))

    # Screenshot area explanation
    add_text(slide, 0.3, 1.5, 9, 0.4,
             "Field 5: Evidence 1", size=9, bold=True, color=RGBColor(0,0,200))
    add_text(slide, 0.3, 1.9, 9, 0.5,
             "Screenshots or copy/pastes from company-official websites, documentation, APIs, SDKs, "
             "presentations, YouTube videos (URL only), etc. demonstrating evidence of use/infringement.",
             size=9, color=RGBColor(80,80,80))

    add_text(slide, 0.3, 2.6, 9, 0.4,
             "Field 6: Evidence 2", size=9, bold=True, color=RGBColor(0,0,200))
    add_text(slide, 0.3, 3.0, 9, 0.4,
             "Additional screenshots. Highlight the specific parts showing infringement.",
             size=9, color=RGBColor(80,80,80))

    add_text(slide, 0.3, 3.6, 9, 0.4,
             "Field 7: Source", size=9, bold=True, color=RGBColor(0,0,200))
    add_text(slide, 0.3, 4.0, 9, 0.3,
             "URL(s) tied to evidence. Each evidence item may have its own URL.",
             size=9, color=RGBColor(80,80,80))

    add_text(slide, 0.3, 4.6, 9, 0.4,
             "Field 8: Student Commentary", size=9, bold=True, color=RGBColor(0,0,200))
    add_text(slide, 0.3, 5.0, 9, 0.8,
             "(1) Why you believe there is infringement\n"
             "(2) How you know: (a) it's the only way it could be done technically, "
             "or (b) we'd need XYZ (source code, etc.) to prove it.",
             size=9, color=RGBColor(80,80,80))

    add_text(slide, 0.3, 6.0, 9, 0.4,
             f"Field 9/10: Company: {company}  |  Product: {product}",
             size=9, color=RGBColor(0,0,200))
    add_footer(slide, slide_num)
    slide_num += 1

    # ---- Element slides ----
    for elem in elements:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Claim limitation
        add_text(slide, 0.3, 0.2, 9.4, 0.6,
                 elem["claim_text"],
                 size=12, bold=True)

        add_text(slide, 0.3, 0.75, 3, 0.2,
                 "Field 4: Claim Limitation",
                 size=8, color=RGBColor(0,0,200))

        # Screenshot placeholder with instructions
        add_screenshot_placeholder(slide, 0.3, 1.1, 4.3, 2.2,
                                   elem["screenshot_label"])

        # Evidence description next to screenshot area
        add_text(slide, 4.8, 1.0, 4.9, 0.2,
                 "Evidence Description:", size=8, bold=True, color=RGBColor(0,0,200))
        ev_box = add_text(slide, 4.8, 1.25, 4.9, 2.0,
                          elem["evidence_desc"], size=9)

        # Source URLs
        add_text(slide, 0.3, 3.5, 2, 0.2,
                 "Field 7: Source URLs", size=8, bold=True, color=RGBColor(0,0,200))
        y = 3.75
        for src in elem["sources"]:
            add_text(slide, 0.3, y, 9.4, 0.22,
                     src, size=8, color=RGBColor(0, 80, 180))
            y += 0.22

        # Second screenshot area (if applicable)
        add_screenshot_placeholder(slide, 0.3, y + 0.15, 4.3, 1.3,
                                   "[INSERT EVIDENCE 2 SCREENSHOT]")
        add_text(slide, 4.8, y + 0.15, 4.9, 0.2,
                 "Additional Evidence / Highlight:", size=8, bold=True, color=RGBColor(0,0,200))

        # Commentary
        add_text(slide, 0.3, y + 1.65, 2, 0.2,
                 "Field 8: Commentary", size=8, bold=True, color=RGBColor(0,0,200))
        comm_box = add_text(slide, 0.3, y + 1.9, 9.4, 1.2,
                            elem["commentary"], size=9)
        comm_box.fill.solid()
        comm_box.fill.fore_color.rgb = RGBColor(255, 255, 235)

        add_footer(slide, slide_num)
        slide_num += 1

    filepath = os.path.join(OUTPUT_DIR, filename)
    prs.save(filepath)
    print(f"Created: {filepath} ({slide_num - 1} slides)")


# =====================================================================
# TARGET DATA
# =====================================================================

google_elements = [
    {
        "claim_text": "(a) obtaining a 2D image of an object;",
        "evidence_desc": (
            "Google's 3D product viewer accepts as few as 1-3 product images. Their research blog states: "
            "'as few as three images capturing most object surfaces were sufficient to improve quality.' "
            "The Veo-based Gen 3 approach can work from a single image."
        ),
        "screenshot_label": "[SCREENSHOT: Google research blog showing\n'as few as three images' input requirement]",
        "sources": [
            "https://research.google/blog/bringing-3d-shoppable-products-online-with-generative-ai/",
        ],
        "commentary": (
            "Google's system clearly obtains 2D images as the starting input. Their own blog documents "
            "accepting 1-3 images per product. This is the foundational input step. "
            "Screenshot the section of the blog where input requirements are described."
        ),
    },
    {
        "claim_text": "(b) classifying the object in the image;",
        "evidence_desc": (
            "Google's system handles diverse product categories: shoes, furniture, apparel, electronics. "
            "The Veo model was fine-tuned on 'millions of high-quality 3D synthetic assets' organized by category. "
            "Google Merchant Center requires product category selection (taxonomy). "
            "The system applies different processing per category."
        ),
        "screenshot_label": "[SCREENSHOT: Google Merchant Center product\ncategory taxonomy / category selection UI]",
        "sources": [
            "https://research.google/blog/bringing-3d-shoppable-products-online-with-generative-ai/",
            "https://support.google.com/merchants/answer/6324436",
        ],
        "commentary": (
            "Google Merchant Center requires explicit product categorization. The 3D generation pipeline "
            "was trained on assets organized by category, implying classification is part of the process. "
            "Screenshot: (1) the blog section mentioning category-specific training, (2) Google Merchant "
            "Center product taxonomy page showing category requirements."
        ),
    },
    {
        "claim_text": "(c) segmenting the object from the image;",
        "evidence_desc": (
            "Google's research blog explicitly confirms 'removing unwanted backgrounds' as a sub-problem "
            "addressed in their pipeline. Gen 1 used 'object-centric images' as input. The product must "
            "be isolated from its background before 3D generation."
        ),
        "screenshot_label": "[SCREENSHOT: Google blog section mentioning\n'removing unwanted backgrounds']",
        "sources": [
            "https://research.google/blog/bringing-3d-shoppable-products-online-with-generative-ai/",
        ],
        "commentary": (
            "Background removal IS segmentation — isolating the object from the scene. Google's own "
            "documentation confirms this step. This directly and unambiguously satisfies element (c). "
            "Screenshot the specific paragraph mentioning background removal."
        ),
    },
    {
        "claim_text": "(d) dimensionally sampling the segmented image of the object;",
        "evidence_desc": (
            "Gen 1: 'predicting 3D priors' and 'estimating camera positions from a sparse set of "
            "object-centric images.' Gen 2: Score Distillation Sampling where 'the 3D model is rendered "
            "from a random camera view, then a view-conditioned diffusion model generates a target from "
            "the same view.' Gen 3: Veo fine-tuned on assets 'rendered from various camera angles.'"
        ),
        "screenshot_label": "[SCREENSHOT: Blog diagram showing\n3D prior extraction / camera pose estimation]",
        "sources": [
            "https://research.google/blog/bringing-3d-shoppable-products-online-with-generative-ai/",
        ],
        "commentary": (
            "All three generations extract dimensional/spatial information from the segmented images — "
            "3D priors, camera positions, depth, and spatial relationships. This is dimensional sampling. "
            "It is technically impossible to generate a 3D model from a 2D image without some form of "
            "dimensional sampling. Screenshot the pipeline diagrams from the blog."
        ),
    },
    {
        "claim_text": "(e) extracting texture information from the segmented image;",
        "evidence_desc": (
            "Google's blog explicitly states Veo captures 'complex interactions between light, material, "
            "texture, and geometry' and handles 'complex lighting and material interactions (i.e., shiny "
            "surfaces).' The output preserves photorealistic product textures, colors, and material properties."
        ),
        "screenshot_label": "[SCREENSHOT: Blog section on Veo's\n'light, material, texture, and geometry' handling]",
        "sources": [
            "https://research.google/blog/bringing-3d-shoppable-products-online-with-generative-ai/",
        ],
        "commentary": (
            "Google explicitly names texture as a core capability of their pipeline. The word 'texture' "
            "appears in their own technical description. This is a direct match. "
            "Screenshot the paragraph mentioning texture and material handling."
        ),
    },
    {
        "claim_text": "(f) generating a 3D mesh model based at least on the dimensional sampling;",
        "evidence_desc": (
            "Gen 1 used NeRF to generate 3D representations. Gen 2 optimized '3D model parameters' via "
            "Score Distillation Sampling. Google also accepts merchant-uploaded glTF/glb files (mesh formats). "
            "Google support docs distinguish between 'glTF/glb models provided by merchants' and "
            "'360-spins generated by Google AI.'"
        ),
        "screenshot_label": "[SCREENSHOT: Google Merchant Center 3D model\nupload page showing glTF/glb format support]",
        "sources": [
            "https://research.google/blog/bringing-3d-shoppable-products-online-with-generative-ai/",
            "https://support.google.com/merchants/answer/13675100",
        ],
        "commentary": (
            "Google generates 3D representations from spatial information extracted from images. Whether "
            "the intermediate representation is NeRF, mesh, or volumetric, the functional result is a 3D model. "
            "The glTF/glb ecosystem confirms mesh-based output is part of their platform. "
            "Screenshot: (1) blog pipeline diagrams, (2) Merchant Center 3D upload requirements."
        ),
    },
    {
        "claim_text": "(g) rendering the texture information onto the 3D mesh model to create a 3D representation of the object.",
        "evidence_desc": (
            "The final output is interactive, photorealistic 360° product visualizations on Google Shopping. "
            "Products display realistic texture, lighting, and material properties. Google reports '50% more "
            "user engagement' for products with 3D imagery. Users rotate and interact with the textured model."
        ),
        "screenshot_label": "[SCREENSHOT: Google Shopping product page\nshowing interactive 3D product viewer]",
        "sources": [
            "https://www.google.com/shopping",
            "https://support.google.com/merchants/answer/13671720",
        ],
        "commentary": (
            "The rendered output is a textured 3D representation viewable from any angle. This is the "
            "final step of the claimed method. Also satisfies dependent Claims 9 (multiple views), "
            "10 (user-input manipulation), and 11 (2D capture of each view as video frames). "
            "Screenshot an actual Google Shopping product with the 3D viewer active."
        ),
    },
]

meshy_elements = [
    {
        "claim_text": "(a) obtaining a 2D image of an object;",
        "evidence_desc": (
            "Meshy's primary workflow: 'Upload a clear image or type a short text prompt.' "
            "The Image-to-3D feature accepts a single image upload as the input to generate a 3D model. "
            "Standard image formats supported."
        ),
        "screenshot_label": "[SCREENSHOT: Meshy.ai upload interface\nshowing single image input]",
        "sources": [
            "https://www.meshy.ai/",
            "https://www.meshy.ai/features/image-to-3d",
        ],
        "commentary": (
            "Meshy's core product offering is single-image to 3D. The upload UI clearly shows image input. "
            "Screenshot the main Image-to-3D page showing the upload interface."
        ),
    },
    {
        "claim_text": "(b) classifying the object in the image;",
        "evidence_desc": (
            "Meshy generates coherent multi-view predictions (front, side, back) from a single input image. "
            "This requires understanding the object type — a shoe, character, vehicle, etc. — to predict "
            "unseen views correctly. The system handles diverse categories differently."
        ),
        "screenshot_label": "[SCREENSHOT: Meshy showing multi-view\ngeneration from single image input]",
        "sources": [
            "https://www.meshy.ai/",
        ],
        "commentary": (
            "Multi-view prediction inherently requires object classification. Generating a correct back "
            "view of a shoe vs. a chair requires different learned priors per object class. While Meshy "
            "may not use an explicit 'classifier' label, the neural network must internally classify "
            "the object to produce coherent views. Source code review would confirm architecture details."
        ),
    },
    {
        "claim_text": "(c) segmenting the object from the image;",
        "evidence_desc": (
            "Meshy isolates the object from the background to generate a clean 3D model. "
            "The output 3D models contain only the target object, not background elements. "
            "Background removal/segmentation is a necessary preprocessing step."
        ),
        "screenshot_label": "[SCREENSHOT: Meshy output showing isolated\n3D object with no background]",
        "sources": [
            "https://www.meshy.ai/",
        ],
        "commentary": (
            "The output proves segmentation occurred — the 3D model contains only the target object, "
            "demonstrating the background was removed. It is technically impossible to generate a clean "
            "3D mesh from a cluttered image without segmentation. Screenshot a Meshy output alongside "
            "its input image to show the background was removed."
        ),
    },
    {
        "claim_text": "(d) dimensionally sampling the segmented image of the object;",
        "evidence_desc": (
            "Meshy's Smart Remesh feature allows adjustment of 'triangle or quad counts' from '1k to 300k.' "
            "The system extracts geometric/dimensional information from the input image to construct "
            "3D vertex positions. Different mesh resolutions represent different sampling granularities."
        ),
        "screenshot_label": "[SCREENSHOT: Meshy Smart Remesh UI showing\ntriangle/quad count adjustment (1k-300k)]",
        "sources": [
            "https://www.meshy.ai/",
        ],
        "commentary": (
            "The adjustable polygon count (1k-300k) directly demonstrates dimensional sampling at "
            "different granularities. The system must sample spatial information from the 2D image to "
            "construct 3D vertices. Screenshot the Smart Remesh interface."
        ),
    },
    {
        "claim_text": "(e) extracting texture information from the segmented image;",
        "evidence_desc": (
            "Meshy generates full PBR (Physically Based Rendering) texture maps: "
            "'Diffuse, Roughness, Metallic, and Normal maps.' Multiple texture styles: "
            "'realistic and cartoon to hand painted and fantasy.' "
            "Separate AI Texturing feature allows re-texturing existing models."
        ),
        "screenshot_label": "[SCREENSHOT: Meshy PBR texture output showing\nDiffuse, Roughness, Metallic, Normal maps]",
        "sources": [
            "https://www.meshy.ai/",
            "https://www.meshy.ai/features/ai-texturing",
        ],
        "commentary": (
            "Meshy's PBR texture pipeline is comprehensive and explicitly documented. Diffuse, Roughness, "
            "Metallic, and Normal maps are standard texture extraction outputs. This is the strongest "
            "possible match for element (e) — no ambiguity. Screenshot the texture output panel."
        ),
    },
    {
        "claim_text": "(f) generating a 3D mesh model based at least on the dimensional sampling;",
        "evidence_desc": (
            "Meshy generates explicit 3D meshes exportable in 7 formats: "
            "'FBX, GLB, OBJ, STL, 3MF, USDZ, BLEND.' "
            "Smart Remesh produces meshes with configurable polygon counts. "
            "These are all standard polygon mesh formats (vertices, faces, edges)."
        ),
        "screenshot_label": "[SCREENSHOT: Meshy export dialog showing\nFBX, GLB, OBJ, STL format options]",
        "sources": [
            "https://www.meshy.ai/",
        ],
        "commentary": (
            "This is the cleanest possible match for element (f). OBJ, FBX, GLB, STL are unambiguously "
            "mesh formats. There is zero design-around argument here — Meshy outputs a '3D mesh model' "
            "by definition. Screenshot the export format selection dialog."
        ),
    },
    {
        "claim_text": "(g) rendering the texture information onto the 3D mesh model to create a 3D representation of the object.",
        "evidence_desc": (
            "The final output is a textured 3D mesh viewable in Meshy's browser-based 3D viewer. "
            "PBR textures are applied to the mesh geometry. Users rotate, zoom, and inspect from any angle. "
            "Models can be downloaded with textures baked onto the mesh."
        ),
        "screenshot_label": "[SCREENSHOT: Meshy 3D viewer showing\ntextured model with rotation controls]",
        "sources": [
            "https://www.meshy.ai/",
        ],
        "commentary": (
            "The browser viewer shows the final textured 3D representation. This also satisfies "
            "Claims 9 (multiple views), 10 (user rotation), and 11 (export of views). "
            "Screenshot a completed model in the browser viewer from multiple angles."
        ),
    },
]

apple_elements = [
    {
        "claim_text": "(a) obtaining a 2D image of an object;",
        "evidence_desc": (
            "SHARP takes a single photograph as input. GitHub README: 'sharp predict -i image_path.png'. "
            "Apple's research page: 'Given a single photograph, it regresses the parameters of a 3D "
            "Gaussian representation of the depicted scene.'"
        ),
        "screenshot_label": "[SCREENSHOT: Apple ML-SHARP GitHub README\nshowing single image CLI command]",
        "sources": [
            "https://github.com/apple/ml-sharp",
            "https://machinelearning.apple.com/research/sharp-monocular-view",
        ],
        "commentary": (
            "SHARP's entire premise is single-image input. The GitHub repo, paper, and Apple ML page "
            "all confirm this. Screenshot the README section showing the CLI command and input description."
        ),
    },
    {
        "claim_text": "(b) classifying the object in the image;",
        "evidence_desc": (
            "SHARP demonstrates 'robust zero-shot generalization across datasets.' The model handles "
            "diverse object and scene categories. Trained on data spanning multiple object types. "
            "The neural network learns category-specific 3D priors during training."
        ),
        "screenshot_label": "[SCREENSHOT: SHARP project page showing\nzero-shot generalization across categories]",
        "sources": [
            "https://apple.github.io/ml-sharp/",
            "https://huggingface.co/papers/2512.10685",
        ],
        "commentary": (
            "Zero-shot generalization implies the model has learned object category representations — "
            "an implicit form of classification. The neural network must internally differentiate object "
            "types to generate appropriate 3D structures. This is a LIKELY match; explicit classification "
            "would need source code analysis to confirm. Screenshot the results showing different categories."
        ),
    },
    {
        "claim_text": "(c) segmenting the object from the image;",
        "evidence_desc": (
            "SHARP uses Apple's Depth Pro for monocular depth estimation. Depth Pro segments the scene "
            "into foreground/background by depth. Per-pixel 3D Gaussian prediction requires understanding "
            "object boundaries. SHARP predicts different Gaussian parameters for object vs. background."
        ),
        "screenshot_label": "[SCREENSHOT: SHARP depth estimation output\nshowing object/background separation]",
        "sources": [
            "https://github.com/apple/ml-sharp",
            "https://github.com/apple/ml-depth-pro",
        ],
        "commentary": (
            "Depth estimation inherently involves segmentation — distinguishing foreground from background "
            "to assign correct depth values. LIKELY match. Explicit segmentation masks may or may not be "
            "an intermediate output. Screenshot depth maps from SHARP showing object isolation."
        ),
    },
    {
        "claim_text": "(d) dimensionally sampling the segmented image of the object;",
        "evidence_desc": (
            "SHARP predicts 'position and appearance of 3D Gaussians' — each with 3D position, scale, "
            "rotation, and opacity. The output is metric-scale: 'absolute scale, supporting metric camera "
            "movements.' Depth Pro provides per-pixel depth (dimensional) information."
        ),
        "screenshot_label": "[SCREENSHOT: SHARP paper figure showing\nGaussian position/scale prediction]",
        "sources": [
            "https://machinelearning.apple.com/research/sharp-monocular-view",
            "https://huggingface.co/papers/2512.10685",
        ],
        "commentary": (
            "Predicting metric-scale 3D positions from a 2D image IS dimensional sampling. Each Gaussian's "
            "position, scale, and orientation directly encode dimensional information extracted from the image. "
            "Screenshot the paper figures showing the prediction pipeline."
        ),
    },
    {
        "claim_text": "(e) extracting texture information from the segmented image;",
        "evidence_desc": (
            "SHARP predicts 'appearance' parameters (spherical harmonic coefficients) for each 3D Gaussian, "
            "encoding color and view-dependent appearance from the input. The output achieves 'photorealistic' "
            "quality with 'sharp details and fine structures.'"
        ),
        "screenshot_label": "[SCREENSHOT: SHARP output examples showing\nphotorealistic texture preservation]",
        "sources": [
            "https://apple.github.io/ml-sharp/",
        ],
        "commentary": (
            "Spherical harmonic appearance parameters encode texture/color information from the input image. "
            "This is functionally equivalent to texture extraction. The photorealistic output quality "
            "demonstrates high-fidelity texture capture. Screenshot example outputs from the project page."
        ),
    },
    {
        "claim_text": "(f) generating a 3D mesh model based at least on the dimensional sampling;",
        "evidence_desc": (
            "SHARP outputs a 3D Gaussian Splatting representation (.ply format). This is a structured 3D "
            "representation with explicit positions, scales, rotations, and spherical harmonic coefficients. "
            "Compatible with standard 3DGS renderers. NOTE: This is NOT a traditional polygon mesh."
        ),
        "screenshot_label": "[SCREENSHOT: SHARP .ply output file and\n3DGS renderer visualization]",
        "sources": [
            "https://github.com/apple/ml-sharp",
            "https://huggingface.co/apple/Sharp",
        ],
        "commentary": (
            "CLAIM CONSTRUCTION ISSUE: SHARP outputs 3D Gaussian splats, not a traditional polygon mesh. "
            "The patent says '3D mesh model.' Arguments FOR: (1) functional equivalence — both are 3D "
            "representations of the object; (2) Gaussian splats can be converted to meshes; (3) Apple's "
            "ecosystem (RealityKit, Object Capture) works with mesh formats. Arguments AGAINST: Gaussians "
            "are mathematically distinct from meshes (no vertices/faces/edges). This element requires "
            "claim construction analysis by patent counsel."
        ),
    },
    {
        "claim_text": "(g) rendering the texture information onto the 3D mesh model to create a 3D representation of the object.",
        "evidence_desc": (
            "SHARP renders photorealistic novel views at '100+ FPS on a standard GPU.' The textured 3D "
            "representation supports 'high-resolution rendering of nearby views with sharp details.' "
            "Integrated into Splat Studio app for Vision Pro. Apple commercialized related tech as "
            "'Spatial Scenes' in iOS 26."
        ),
        "screenshot_label": "[SCREENSHOT: Splat Studio on Vision Pro /\nSpatial Scenes iOS 26 demo]",
        "sources": [
            "https://apple.github.io/ml-sharp/",
            "https://www.uploadvr.com/apple-sharp-open-source-on-device-gaussian-splatting/",
        ],
        "commentary": (
            "The photorealistic rendered output satisfies the functional requirement of element (g). "
            "Vision Pro integration (Splat Studio) and iOS 26 Spatial Scenes demonstrate commercial "
            "deployment. Satisfies Claims 9 (novel views) and 10 (user viewpoint movement). "
            "Screenshot Splat Studio and Spatial Scenes demos."
        ),
    },
]

amazon_elements = [
    {
        "claim_text": "(a) obtaining a 2D image of an object;",
        "evidence_desc": (
            "Amazon Seller App 'Create 3D Models' captures product images via iOS camera. "
            "Amazon requires '2-10 reference photos' for 3D model creation. "
            "Certified provider Nextech3D.ai generates 3D models from 1-2 product photos. "
            "AWS Human Mesh Recovery processes 2D images as input."
        ),
        "screenshot_label": "[SCREENSHOT: Amazon Seller App 'Create 3D\nModels' interface showing image capture]",
        "sources": [
            "https://sell.amazon.com/tools/3d-ar",
            "https://m.media-amazon.com/images/G/01/imaging/3D/Getting_Started_With_3D_Content_on_Amazon_-_Spring_2023_US.pdf",
        ],
        "commentary": (
            "Amazon's 3D pipeline starts with 2D images at every entry point — seller scanning, reference "
            "photo upload, and third-party provider input. Screenshot the Seller App 3D scanning interface "
            "and the Getting Started guide showing photo requirements."
        ),
    },
    {
        "claim_text": "(b) classifying the object in the image;",
        "evidence_desc": (
            "Amazon explicitly classifies products by category for 3D eligibility: 'Eligibility depends "
            "on product category, physical properties, and other factors.' Different 3D experiences per "
            "category: Virtual Try-On (shoes/eyewear), View in Your Room (furniture), View in 3D (general). "
            "Amazon product taxonomy assigns every product to a category."
        ),
        "screenshot_label": "[SCREENSHOT: Amazon Seller Central showing\nproduct category / 3D eligibility criteria]",
        "sources": [
            "https://sell.amazon.com/tools/3d-ar",
            "https://sellercentral.amazon.com/",
        ],
        "commentary": (
            "Amazon's category-dependent 3D processing is explicit classification. Products must be "
            "categorized before 3D generation can occur. The system selects different 3D experiences "
            "(Try-On, View in Room, View in 3D) based on category. Screenshot the seller documentation "
            "showing category eligibility for 3D features."
        ),
    },
    {
        "claim_text": "(c) segmenting the object from the image;",
        "evidence_desc": (
            "AWS Spatial Computing blog describes using 'image segmentation and semantic labelling' via "
            "Amazon SageMaker and Amazon Bedrock in their 3D reconstruction pipeline. "
            "The Seller App scanning isolates the product from the scanning environment. "
            "3D models contain only the product, not surroundings."
        ),
        "screenshot_label": "[SCREENSHOT: AWS blog mentioning 'image\nsegmentation and semantic labelling']",
        "sources": [
            "https://aws.amazon.com/blogs/spatial/3d-gaussian-splatting-performant-3d-scene-reconstruction-at-scale/",
        ],
        "commentary": (
            "AWS explicitly documents 'image segmentation' as part of their 3D pipeline. This is direct "
            "evidence from Amazon's own technical blog. Screenshot the specific paragraph mentioning "
            "segmentation in the AWS Spatial Computing blog post."
        ),
    },
    {
        "claim_text": "(d) dimensionally sampling the segmented image of the object;",
        "evidence_desc": (
            "Amazon requires 'accurate product dimensions' submitted with 3D models. The Seller App "
            "scanning extracts spatial data over a 5-10 minute scanning session. AWS Human Mesh Recovery "
            "uses ScoreHMR to 'capture and reconstruct human body parameters from input images.' "
            "AWS 3D Gaussian Splatting pipeline extracts spatial geometry."
        ),
        "screenshot_label": "[SCREENSHOT: Amazon Seller Central 3D model\nsubmission showing dimension requirements]",
        "sources": [
            "https://sell.amazon.com/tools/3d-ar",
            "https://aws.amazon.com/blogs/spatial/from-2d-to-3d-building-a-scalable-human-mesh-recovery-pipeline-with-amazon-sagemaker-ai/",
        ],
        "commentary": (
            "Dimensional sampling occurs during scanning (extracting 3D geometry from 2D captures) and "
            "in the AWS pipeline (extracting body/object parameters). The requirement for 'accurate product "
            "dimensions' confirms dimensional data is core to the process. Screenshot both the seller "
            "dimension requirements and the AWS HMR pipeline diagram."
        ),
    },
    {
        "claim_text": "(e) extracting texture information from the segmented image;",
        "evidence_desc": (
            "Amazon's 3D product models display photorealistic textures — product colors, labels, branding, "
            "material properties are all preserved from the original product. AWS 3D pipeline handles "
            "texture as part of the reconstruction. GLB/GLTF format includes embedded texture data."
        ),
        "screenshot_label": "[SCREENSHOT: Amazon product listing showing\nphotorealistic 3D model with textures]",
        "sources": [
            "https://sell.amazon.com/tools/3d-ar",
            "https://www.amazon.com/",
        ],
        "commentary": (
            "The photorealistic output with preserved labels, colors, and materials proves texture "
            "extraction occurred. GLB/GLTF format embeds texture maps by specification. Screenshot "
            "an actual Amazon product listing with View in 3D active, showing texture fidelity."
        ),
    },
    {
        "claim_text": "(f) generating a 3D mesh model based at least on the dimensional sampling;",
        "evidence_desc": (
            "Amazon requires all 3D models in GLB or GLTF format — both are mesh-based 3D formats "
            "(vertices, faces, materials) by specification. AWS blog describes 'mesh representations.' "
            "Amazon shifted from 360° images to 3D models (mandatory since Dec 14, 2023)."
        ),
        "screenshot_label": "[SCREENSHOT: Amazon 3D content requirements\nshowing GLB/GLTF format mandate]",
        "sources": [
            "https://sell.amazon.com/tools/3d-ar",
            "https://support.google.com/merchants/answer/13675100",
        ],
        "commentary": (
            "GLB and GLTF are mesh formats BY DEFINITION — they contain vertices, faces, normals, and "
            "materials. Amazon mandates these formats for all 3D product content. There is zero ambiguity "
            "that Amazon's pipeline produces '3D mesh models.' Screenshot the Amazon seller documentation "
            "specifying GLB/GLTF format requirements."
        ),
    },
    {
        "claim_text": "(g) rendering the texture information onto the 3D mesh model to create a 3D representation of the object.",
        "evidence_desc": (
            "Amazon's 'View in 3D' lets customers 'rotate and zoom in on a product in 3D, offering a "
            "detailed view from all angles.' 'View in Your Room' renders the textured 3D model in AR. "
            "'Virtual Try-On' renders products on the user. Amazon reports 2X purchase conversion and "
            "20% lower returns from 3D features."
        ),
        "screenshot_label": "[SCREENSHOT: Amazon product page with\nView in 3D / View in Your Room active]",
        "sources": [
            "https://sell.amazon.com/tools/3d-ar",
            "https://www.amazon.com/",
        ],
        "commentary": (
            "The interactive textured 3D viewer IS the rendered output of the claimed method. "
            "Satisfies Claims 9 (multiple views — rotate any angle), 10 (user input — rotate/zoom), "
            "and 11 (2D captures — AR view screenshots). Amazon's scale (millions of listings) and "
            "mandatory 3D adoption make this a high-value infringement target. Screenshot: actual "
            "Amazon product page with 3D viewer open, and View in Your Room AR screenshot."
        ),
    },
]


if __name__ == "__main__":
    build_eou("EoU_Google_Shopping.pptx", "Google (Alphabet Inc.)",
              "Google Shopping 3D Product Visualizations", google_elements)
    build_eou("EoU_Meshy_AI.pptx", "Meshy Inc.",
              "Meshy AI 3D Model Generator", meshy_elements)
    build_eou("EoU_Apple_SHARP.pptx", "Apple Inc.",
              "SHARP + Spatial Scenes (iOS 26) + Vision Pro", apple_elements)
    build_eou("EoU_Amazon.pptx", "Amazon.com, Inc.",
              "Amazon 3D/AR (Seller App + View in 3D + AWS Pipeline)", amazon_elements)
    print("\nAll 4 EoU decks created.")
